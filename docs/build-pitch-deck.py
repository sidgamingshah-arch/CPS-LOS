"""
Helix pitch deck builder — CRISIL brand palette.

Produces docs/Helix-Pitch-Deck.pptx, 16:9 widescreen, ~22 slides.

Design language (CRISIL):
- Brand voltage: Ocean 8 #005B72 (dominant).
- Surface system: white canvas, deep-ocean dark hero (#042e39), Light-Blue
  (#C7DAE0) tinted panels with Ocean numerals, Warm-Gray (#C2C4C3) card borders.
- Semantic: positive #1B8D2F, negative #B20023.
- Governance: AI=purple #782080, Human=green #1B8D2F, Deterministic=Ocean (semantic).
- Type: Inter (display + body), JetBrains Mono via "Consolas" fallback for figures.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

# ─── CRISIL palette ─────────────────────────────────────────────────────────
OCEAN_8        = RGBColor(0x00, 0x5B, 0x72)
OCEAN_7        = RGBColor(0x00, 0x49, 0x5B)
OCEAN_6        = RGBColor(0x00, 0x6D, 0x89)
OCEAN_5        = RGBColor(0x49, 0x97, 0xAB)
OCEAN_4        = RGBColor(0x6D, 0xAC, 0xBC)
LIGHT_BLUE     = RGBColor(0xC7, 0xDA, 0xE0)
LIGHT_BLUE_SOFT= RGBColor(0xE9, 0xF1, 0xF4)
WARM_GRAY      = RGBColor(0xC2, 0xC4, 0xC3)
HAIRLINE       = RGBColor(0xD7, 0xDC, 0xDD)
HAIRLINE_SOFT  = RGBColor(0xE7, 0xEA, 0xEA)
INK            = RGBColor(0x0C, 0x14, 0x17)
BODY           = RGBColor(0x44, 0x4D, 0x52)
MUTED          = RGBColor(0x6B, 0x73, 0x78)
MUTED_SOFT     = RGBColor(0x9A, 0xA1, 0xA5)
CANVAS         = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE_SOFT   = RGBColor(0xF5, 0xF7, 0xF8)
SURFACE_DARK   = RGBColor(0x04, 0x2E, 0x39)
SURFACE_DARK_2 = RGBColor(0x0A, 0x41, 0x50)

# Governance (semantic)
AI_PURPLE      = RGBColor(0x78, 0x20, 0x80)
AI_PURPLE_SOFT = RGBColor(0xF1, 0xE6, 0xF2)
HUMAN_GREEN    = RGBColor(0x1B, 0x8D, 0x2F)
HUMAN_GREEN_SFT= RGBColor(0xE3, 0xF3, 0xE6)
DET_OCEAN      = OCEAN_7
DET_OCEAN_SFT  = LIGHT_BLUE_SOFT

# Semantic
POSITIVE       = RGBColor(0x1B, 0x8D, 0x2F)
NEGATIVE       = RGBColor(0xB2, 0x00, 0x23)
ACCENT_YELLOW  = RGBColor(0xED, 0xF8, 0x6F)

# ─── Deck geometry ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ─── Type ───────────────────────────────────────────────────────────────────
FONT_SANS = "Inter"
FONT_MONO = "Consolas"  # JetBrains Mono substitute available in PPT

ASSETS = "/tmp/helix-pitch"

# ─── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = line_w
    if not shadow:
        _disable_shadow(shp)
    return shp

def add_round_rect(slide, x, y, w, h, fill, radius_frac=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = line_w
    if radius_frac is not None:
        # set the corner radius (adjustment 0..0.5 of the smaller side)
        shp.adjustments[0] = radius_frac
    _disable_shadow(shp)
    return shp

def _disable_shadow(shape):
    # Remove the default soft shadow that PPT adds to shapes.
    sp = shape._element
    spPr = sp.find('.//' + qn('p:spPr'))
    if spPr is None: return
    # remove existing effect lists
    for tag in ('a:effectLst', 'a:effectDag'):
        for n in spPr.findall(qn(tag)):
            spPr.remove(n)
    # explicit empty effectLst disables theme-default shadows
    spPr.append(etree.SubElement(spPr, qn('a:effectLst')))

def add_text(slide, x, y, w, h, text, *,
             font=FONT_SANS, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, italic=False, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    # Use first paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        # spacing in 100ths of a point (positive = wider tracking)
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(letter_spacing * 100)))
    return tb

def add_multi(slide, x, y, w, h, paras, *,
              font=FONT_SANS, anchor=MSO_ANCHOR.TOP):
    """paras: list of (text, size, color, bold, line_spacing, align)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(paras):
        text, size, color, bold, ls, align = (item + (None,)*6)[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align: p.alignment = align
        if ls: p.line_spacing = ls
        run = p.add_run(); run.text = text
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bool(bold)
        if color: run.font.color.rgb = color
    return tb

def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        # Placeholder if image is missing
        add_rect(slide, x, y, w or Inches(4), h or Inches(2.5), SURFACE_SOFT, line=HAIRLINE)
        add_text(slide, x, y, w or Inches(4), h or Inches(2.5), "[image not found: " + os.path.basename(path) + "]",
                 color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, size=10)
        return None
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

# Pill / chip / badge primitives
def add_pill(slide, x, y, w, h, text, fill, text_color, *, size=8, font=FONT_SANS, bold=True, letter_spacing=0.6):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.adjustments[0] = 0.5
    _disable_shadow(shp)
    tf = shp.text_frame
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = text_color
    rPr = run._r.get_or_add_rPr(); rPr.set('spc', str(int(letter_spacing * 100)))
    return shp

def add_gov_chips(slide, x, y, kinds=("ai", "human", "det"), small=False):
    """Render the AI · ADVISORY / HUMAN-GATED / DETERMINISTIC pill cluster."""
    chips = {
        "ai":    ("● AI · ADVISORY", AI_PURPLE_SOFT, AI_PURPLE),
        "human": ("● HUMAN-GATED",   HUMAN_GREEN_SFT, HUMAN_GREEN),
        "det":   ("● DETERMINISTIC", DET_OCEAN_SFT, DET_OCEAN),
    }
    h = Inches(0.28) if small else Inches(0.34)
    sz = 8 if small else 9
    gap = Inches(0.10)
    cx = x
    pills = []
    for k in kinds:
        label, fill, color = chips[k]
        # measure approx width by char count
        w = Inches(0.10 + 0.075 * len(label.replace("● ", "")) )
        p = add_pill(slide, cx, y, w, h, label, fill, color, size=sz)
        pills.append(p)
        cx += w + gap
    return cx

def add_hairline(slide, x, y, w, color=HAIRLINE, weight=Pt(0.75)):
    line = slide.shapes.add_connector(1, x, y, x + w, y)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = weight
    return line

def add_topbar(slide, page_title, crumb=None, *, gov=True, page_num=None, page_of=None):
    """The thin topbar with crumb + title + governance chips."""
    H = Inches(0.95)
    add_rect(slide, 0, 0, SW, H, CANVAS)
    add_hairline(slide, 0, H, SW)
    if crumb:
        add_text(slide, Inches(0.6), Inches(0.18), Inches(7), Inches(0.22),
                 crumb, size=10, color=MUTED)
    add_text(slide, Inches(0.6), Inches(0.40), Inches(8), Inches(0.5),
             page_title, size=22, bold=True, color=INK, letter_spacing=-0.4)
    if gov:
        add_gov_chips(slide, Inches(9.8), Inches(0.32))
    if page_num and page_of:
        add_text(slide, Inches(12.4), Inches(0.32), Inches(0.7), Inches(0.3),
                 f"{page_num} / {page_of}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    return H

def add_footer(slide, *, version="v1.0"):
    y = SH - Inches(0.35)
    add_text(slide, Inches(0.6), y, Inches(7), Inches(0.25),
             "Helix · Governed AI for Wholesale Credit · CRISIL",
             size=9, color=MUTED)
    add_text(slide, Inches(11), y, Inches(2), Inches(0.25),
             "Confidential · " + version, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── Slide builders ─────────────────────────────────────────────────────────

PAGE_OF = 22  # filled at the end after the deck is fixed

def s_title():
    s = prs.slides.add_slide(BLANK)
    # full-bleed deep-ocean editorial canvas
    add_rect(s, 0, 0, SW, SH, SURFACE_DARK)
    # subtle accent shapes (ocean ramp)
    accent = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2.5), Inches(8), Inches(8))
    accent.fill.solid(); accent.fill.fore_color.rgb = OCEAN_6
    accent.line.fill.background()
    _disable_shadow(accent)
    accent.element.find('.//' + qn('p:spPr'))  # ensure spPr
    # adjust transparency on the accent oval via XML
    sppr = accent.fill._xPr.find(qn('a:solidFill'))
    sppr.find(qn('a:srgbClr')).append(etree.fromstring('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="14000"/>'))

    accent2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(3.5), Inches(8), Inches(6))
    accent2.fill.solid(); accent2.fill.fore_color.rgb = AI_PURPLE
    accent2.line.fill.background()
    _disable_shadow(accent2)
    sppr = accent2.fill._xPr.find(qn('a:solidFill'))
    sppr.find(qn('a:srgbClr')).append(etree.fromstring('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="10000"/>'))

    # Top brand strip
    add_text(s, Inches(0.85), Inches(0.6), Inches(6), Inches(0.4),
             "CRISIL", size=11, bold=True, color=CANVAS, letter_spacing=2)

    # Big wordmark
    add_text(s, Inches(0.85), Inches(2.2), Inches(7), Inches(1.6),
             "Helix", size=110, bold=False, color=CANVAS, letter_spacing=-3, line_spacing=1.0)
    # blue 'x' overlay — simulate by overlaying the same word with just the X colored
    # PPT doesn't support per-character colour easily, so render an accent dot beside it
    add_text(s, Inches(4.7), Inches(2.2), Inches(0.8), Inches(1.6),
             "·", size=110, color=OCEAN_5, line_spacing=1.0)

    add_text(s, Inches(0.85), Inches(4.10), Inches(11), Inches(0.6),
             "Governed AI for Wholesale Credit",
             size=30, color=CANVAS, letter_spacing=-0.6)

    add_text(s, Inches(0.85), Inches(4.95), Inches(11), Inches(0.4),
             "AI where it helps   ·   Humans where regulation demands   ·   Deterministic figures throughout",
             size=14, color=RGBColor(0x9F, 0xB8, 0xC0), letter_spacing=0.4)

    # footer
    add_text(s, Inches(0.85), Inches(6.85), Inches(6), Inches(0.3),
             "Wholesale loan origination · lifecycle · audit · governance",
             size=10, color=RGBColor(0x9F, 0xB8, 0xC0))
    add_text(s, Inches(11), Inches(6.85), Inches(2), Inches(0.3),
             "Pitch deck · v1.0", size=10, color=RGBColor(0x9F, 0xB8, 0xC0), align=PP_ALIGN.RIGHT)

def s_section_break(num, title, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, SURFACE_DARK)
    add_text(s, Inches(0.85), Inches(2.3), Inches(2), Inches(0.5),
             f"{num:02d}", size=22, color=OCEAN_5, font=FONT_MONO, letter_spacing=2)
    add_text(s, Inches(0.85), Inches(2.9), Inches(12), Inches(1.2),
             title, size=44, color=CANVAS, letter_spacing=-1, line_spacing=1.05)
    if sub:
        add_text(s, Inches(0.85), Inches(4.3), Inches(11.5), Inches(0.8),
                 sub, size=15, color=RGBColor(0x9F, 0xB8, 0xC0), line_spacing=1.45)
    # ornament line
    line = s.shapes.add_connector(1, Inches(0.85), Inches(5.6), Inches(2.5), Inches(5.6))
    line.line.color.rgb = OCEAN_5; line.line.width = Pt(2)

def s_problem():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "The problem", crumb="The wholesale-credit AI dilemma", page_num=2, page_of=PAGE_OF)
    # left/right split
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
             "Banks want AI. Regulators want accountability.",
             size=30, bold=False, color=INK, letter_spacing=-0.6, line_spacing=1.15)
    add_text(s, Inches(0.6), Inches(2.05), Inches(12), Inches(0.6),
             "Most platforms force a choice.",
             size=30, bold=False, color=INK, letter_spacing=-0.6, line_spacing=1.15)

    # Two columns: "AI-first tools" vs "Traditional LOS"
    col_y = Inches(3.3)
    col_h = Inches(3.4)
    col_w = Inches(6.0)
    gap   = Inches(0.4)
    # Card 1
    x1 = Inches(0.6)
    add_round_rect(s, x1, col_y, col_w, col_h, CANVAS, radius_frac=0.04, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x1+Inches(0.3), col_y+Inches(0.3), col_w-Inches(0.6), Inches(0.4),
             "AI-FIRST TOOLS", size=10, bold=True, color=AI_PURPLE, letter_spacing=1.2)
    add_text(s, x1+Inches(0.3), col_y+Inches(0.65), col_w-Inches(0.6), Inches(0.7),
             "Fast. Suggestive. Useful.", size=19, bold=False, color=INK, letter_spacing=-0.3)
    bullets1 = [
        "But: no named accountable human on the figure path.",
        "But: AI mutates the credit decision opaquely.",
        "But: no immutable audit a regulator can read.",
    ]
    by = col_y + Inches(1.55)
    for b in bullets1:
        add_text(s, x1+Inches(0.3), by, Inches(0.25), Inches(0.35), "—", size=13, color=NEGATIVE, bold=True)
        add_text(s, x1+Inches(0.55), by, col_w-Inches(0.85), Inches(0.4), b, size=12, color=BODY, line_spacing=1.4)
        by += Inches(0.45)

    # Card 2
    x2 = x1 + col_w + gap
    add_round_rect(s, x2, col_y, col_w, col_h, CANVAS, radius_frac=0.04, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x2+Inches(0.3), col_y+Inches(0.3), col_w-Inches(0.6), Inches(0.4),
             "TRADITIONAL LOS", size=10, bold=True, color=DET_OCEAN, letter_spacing=1.2)
    add_text(s, x2+Inches(0.3), col_y+Inches(0.65), col_w-Inches(0.6), Inches(0.7),
             "Auditable. Deterministic. Slow.", size=19, bold=False, color=INK, letter_spacing=-0.3)
    bullets2 = [
        "But: no AI leverage on extraction, drafting, scoring.",
        "But: every regulation forks the codebase.",
        "But: analyst hours wasted on AI-tractable work.",
    ]
    by = col_y + Inches(1.55)
    for b in bullets2:
        add_text(s, x2+Inches(0.3), by, Inches(0.25), Inches(0.35), "—", size=13, color=NEGATIVE, bold=True)
        add_text(s, x2+Inches(0.55), by, col_w-Inches(0.85), Inches(0.4), b, size=12, color=BODY, line_spacing=1.4)
        by += Inches(0.45)
    add_footer(s)

def s_thesis():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "The thesis", crumb="Helix doesn't force the choice", page_num=3, page_of=PAGE_OF)
    # Centerpiece: three pillars (purple / green / ocean)
    pillars = [
        ("AI", "where it helps",
         "Extract. Classify. Spread. Score. Optimise. Draft. Translate. Suggest.",
         AI_PURPLE, AI_PURPLE_SOFT),
        ("Humans", "where regulation demands",
         "Confirm. Override-with-reason. Approve. Sign. Waive. Decide. SoD-enforced.",
         HUMAN_GREEN, HUMAN_GREEN_SFT),
        ("Deterministic", "figures throughout",
         "Rating. Capital. PD/LGD/EAD. RAROC. ECL. Pricing of record. Byte-identical.",
         DET_OCEAN, DET_OCEAN_SFT),
    ]
    pw = Inches(4.0); gap = Inches(0.2); top = Inches(1.8); ph = Inches(4.4)
    total = pw*3 + gap*2
    start = (SW - total) / 2
    for i, (head, sub, body, color, soft) in enumerate(pillars):
        x = start + i*(pw+gap)
        add_round_rect(s, x, top, pw, ph, soft, radius_frac=0.04, line=color, line_w=Pt(1.0))
        add_pill(s, x+Inches(0.3), top+Inches(0.3), Inches(0.55), Inches(0.28),
                 f"0{i+1}", color, CANVAS, size=8)
        add_text(s, x+Inches(0.3), top+Inches(0.8), pw-Inches(0.6), Inches(0.7),
                 head, size=34, bold=False, color=color, letter_spacing=-0.6, line_spacing=1.0)
        add_text(s, x+Inches(0.3), top+Inches(1.55), pw-Inches(0.6), Inches(0.4),
                 sub, size=13, color=INK, line_spacing=1.3)
        add_text(s, x+Inches(0.3), top+Inches(2.20), pw-Inches(0.6), Inches(1.9),
                 body, size=12, color=BODY, line_spacing=1.55)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.5),
             "Not three features. One operating doctrine — wired into every screen, every API, every audit row.",
             size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s)

def s_what_is_helix():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "What is Helix?", crumb="One sentence, one diagram, one promise", page_num=4, page_of=PAGE_OF)
    # Big one-liner
    add_text(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(1.5),
             "An AI-first wholesale loan origination & lifecycle platform that takes a deal from "
             "prospect to portfolio without ever letting AI move the figure of record.",
             size=22, color=INK, letter_spacing=-0.3, line_spacing=1.35)

    # 6 numbers row
    stats = [
        ("9", "microservices"),
        ("22", "master types"),
        ("334", "e2e safety assertions"),
        ("1", "named accountable human per action"),
        ("0", "AI writes to authoritative figures"),
        ("100%", "actions stamped on an immutable audit"),
    ]
    y = Inches(3.6); h = Inches(1.6)
    n = len(stats); gap = Inches(0.16); avail = SW - Inches(1.2) - gap*(n-1)
    sw = avail / n
    x = Inches(0.6)
    for v, lbl in stats:
        add_round_rect(s, x, y, sw, h, LIGHT_BLUE_SOFT, radius_frac=0.06, line=HAIRLINE, line_w=Pt(0.5))
        add_text(s, x, y+Inches(0.22), sw, Inches(0.55),
                 v, size=30, bold=True, color=OCEAN_8, font=FONT_MONO, align=PP_ALIGN.CENTER, letter_spacing=-0.4)
        add_text(s, x+Inches(0.1), y+Inches(0.95), sw-Inches(0.2), Inches(0.6),
                 lbl, size=10, color=OCEAN_7, align=PP_ALIGN.CENTER, line_spacing=1.3)
        x += sw + gap

    # Below: three sub-paragraphs
    para_y = Inches(5.6)
    cw = Inches(4.0); cg = Inches(0.2)
    blurbs = [
        ("Built around the lifecycle",
         "Counterparty → Origination → Risk → Decision → Limits → Portfolio. One workspace, one audit, one story."),
        ("Governed by construction",
         "Every advisory output is a separate entity, an AI audit stamp, and a human gate. Asserted in e2e."),
        ("Regime as data, not code",
         "Capital, PD/LGD, ECL, DoA, pricing — versioned rule packs the engines read at runtime."),
    ]
    cx = Inches(0.6)
    for title, body in blurbs:
        add_text(s, cx, para_y, cw, Inches(0.4), title, size=12.5, bold=True, color=OCEAN_8)
        add_text(s, cx, para_y+Inches(0.4), cw, Inches(1.4), body, size=11, color=BODY, line_spacing=1.45)
        cx += cw + cg
    add_footer(s)

def s_lifecycle():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "The lifecycle spine", crumb="One deal, end-to-end, one audit", page_num=5, page_of=PAGE_OF)
    stages = [
        ("Counterparty",  "Onboard · KYC/CDD/UBO · groups"),
        ("Origination",   "Spread · structure · doc-intel"),
        ("Risk",          "Rate · capital · RAROC · overlays"),
        ("Decision",      "DoA · CAD · covenants · docs"),
        ("Limits",        "Tree · utilisation · EOD"),
        ("Portfolio",     "Book · ECL · EWS · exports"),
    ]
    n = len(stages); gap = Inches(0.18); margin = Inches(0.6)
    avail = SW - margin*2 - gap*(n-1)
    sw = avail / n
    y = Inches(2.0); h = Inches(2.6)
    x = margin
    for i, (head, sub) in enumerate(stages):
        # card
        add_round_rect(s, x, y, sw, h, CANVAS, radius_frac=0.06, line=HAIRLINE, line_w=Pt(0.75))
        # number
        add_text(s, x+Inches(0.25), y+Inches(0.25), Inches(0.8), Inches(0.4),
                 f"{i+1:02d}", size=10, font=FONT_MONO, color=OCEAN_5, letter_spacing=1.5)
        # name
        add_text(s, x+Inches(0.25), y+Inches(0.65), sw-Inches(0.5), Inches(0.6),
                 head, size=17, bold=True, color=OCEAN_8, letter_spacing=-0.2, line_spacing=1.15)
        # sub
        add_text(s, x+Inches(0.25), y+Inches(1.30), sw-Inches(0.5), Inches(1.2),
                 sub, size=11, color=BODY, line_spacing=1.45)
        # arrow between
        if i < n-1:
            ax = x + sw + Inches(0.01)
            ay = y + h/2 - Inches(0.13)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.16), Inches(0.25))
            arr.fill.solid(); arr.fill.fore_color.rgb = OCEAN_5
            arr.line.fill.background()
            _disable_shadow(arr)
        x += sw + gap

    # under-the-spine banner
    by = Inches(5.0); bh = Inches(1.5)
    add_round_rect(s, Inches(0.6), by, SW-Inches(1.2), bh, LIGHT_BLUE_SOFT, radius_frac=0.04, line=OCEAN_4, line_w=Pt(0.5))
    add_text(s, Inches(0.9), by+Inches(0.20), Inches(11.5), Inches(0.4),
             "One spine, one audit, one accountable human per step", size=14, bold=True, color=OCEAN_8)
    add_text(s, Inches(0.9), by+Inches(0.65), Inches(11.5), Inches(0.8),
             "Helix-common stamps every write with X-Actor and a HUMAN / AI / SYSTEM type. Re-reading the trail "
             "reconstructs the decision, who took it, what figures it depended on, and which rule pack was in force.",
             size=11.5, color=OCEAN_7, line_spacing=1.5)
    add_footer(s)

def s_dashboard():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Portfolio Dashboard", crumb="One book · stages · concentration · model fit", page_num=6, page_of=PAGE_OF)
    add_image(s, f"{ASSETS}/01-dashboard.png", Inches(0.6), Inches(1.10), w=Inches(8.5))
    # right-side annotations
    x = Inches(9.4); y = Inches(1.2); w = Inches(3.4)
    add_text(s, x, y, w, Inches(0.4), "WHAT YOU SEE", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    bullets = [
        ("Light-Blue stat tiles", "Ocean numerals on a CRISIL-tinted panel — every figure deterministic."),
        ("IFRS 9 / Ind AS 109", "Provision split by stage. Deterministic; AI doesn't enter here."),
        ("Concentration & limits", "Single-name + segment caps; breach surfaced as semantic red."),
        ("Model-fit gate", "Override rate is a model-trust signal — alerts above 25% (PRD §11)."),
    ]
    by = y + Inches(0.55)
    for t, b in bullets:
        add_text(s, x, by, w, Inches(0.32), t, size=11.5, bold=True, color=INK)
        add_text(s, x, by+Inches(0.32), w, Inches(0.7), b, size=10.5, color=BODY, line_spacing=1.45)
        by += Inches(1.1)
    add_footer(s)

def s_risk_lab_split():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Risk Lab — the signature split", crumb="AI · ADVISORY ↔ AUTHORITATIVE · UNCHANGED", page_num=7, page_of=PAGE_OF)
    add_image(s, f"{ASSETS}/02-risklab-split.png", Inches(0.6), Inches(1.10), w=Inches(8.5))
    x = Inches(9.4); y = Inches(1.2); w = Inches(3.4)
    add_text(s, x, y, w, Inches(0.4), "THE MOTION CENTREPIECE", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    add_text(s, x, y+Inches(0.5), w, Inches(2.0),
             "Statistical RAG band of GREEN 86.2 / 100 on the left.\n\n"
             "Authoritative rating — AA, PD 0.05% — on the right, with a dashed ● UNCHANGED tag.",
             size=12, color=INK, line_spacing=1.55)
    # callout
    cy = y+Inches(3.0)
    add_round_rect(s, x, cy, w, Inches(2.5), SURFACE_DARK, radius_frac=0.06)
    add_text(s, x+Inches(0.25), cy+Inches(0.25), w-Inches(0.5), Inches(0.4),
             "WHAT THIS PROVES", size=10, bold=True, color=OCEAN_5, letter_spacing=1.2)
    add_text(s, x+Inches(0.25), cy+Inches(0.65), w-Inches(0.5), Inches(1.7),
             "The advisory score and the figure of record sit side-by-side. The e2e suite asserts "
             "this rating is byte-identical before and after the AI runs.",
             size=11, color=CANVAS, line_spacing=1.6)
    add_footer(s)

def s_pricing_lab():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Pricing Lab — pricing of record preserved", crumb="Optimiser · concessions · SoD-approved", page_num=8, page_of=PAGE_OF)
    add_image(s, f"{ASSETS}/03-pricinglab.png", Inches(0.6), Inches(1.10), w=Inches(8.5))
    x = Inches(9.4); y = Inches(1.2); w = Inches(3.4)
    add_text(s, x, y, w, Inches(0.4), "GOVERNANCE CHAIN", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    # vertical chain
    chain = [
        ("Target RAROC", "set by the RM"),
        ("Optimiser proposes", "AI · advisory"),
        ("Concession breach", "−522 bps"),
        ("L1 ✓", "human · maker"),
        ("L2 ✓", "human · checker (SoD)"),
        ("● PRESERVED", "pricing of record unchanged"),
    ]
    cy = y + Inches(0.55)
    for label, sub in chain:
        bullet_color = HUMAN_GREEN if "✓" in label or "PRESERVED" in label else (AI_PURPLE if "Optimiser" in label or "advisory" in sub else OCEAN_8)
        add_text(s, x, cy, Inches(0.3), Inches(0.35), "●", size=14, color=bullet_color, bold=True)
        add_text(s, x+Inches(0.30), cy, w-Inches(0.3), Inches(0.30), label, size=12, bold=True, color=INK)
        add_text(s, x+Inches(0.30), cy+Inches(0.30), w-Inches(0.3), Inches(0.30), sub, size=10, color=MUTED, italic=True)
        cy += Inches(0.5)
    add_footer(s)

def _img_caption_slide(title, crumb, img_path, paras, *, num=0):
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, title, crumb=crumb, page_num=num, page_of=PAGE_OF)
    add_image(s, img_path, Inches(0.6), Inches(1.10), w=Inches(8.5))
    x = Inches(9.4); y = Inches(1.2); w = Inches(3.4)
    add_text(s, x, y, w, Inches(0.4), "WHAT YOU SEE", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    by = y + Inches(0.55)
    for t, b in paras:
        add_text(s, x, by, w, Inches(0.32), t, size=11.5, bold=True, color=INK)
        add_text(s, x, by+Inches(0.32), w, Inches(0.7), b, size=10.5, color=BODY, line_spacing=1.45)
        by += Inches(1.1)
    add_footer(s)
    return s

def s_workspace():
    _img_caption_slide(
        "Deal Workspace", "AI-executed, human-gated, end to end",
        f"{ASSETS}/04-workspace.png",
        [
            ("One deal, one workspace", "Intake → spread → rate → capital → price → decide → book — no service-hopping."),
            ("Facilities + sublimits", "Interchangeability groups; cap rollups; renewals copy-from prior cycle."),
            ("Specialised structures", "Group / joint / dual-obligor / syndication / FI ICR — declared, not coded."),
            ("Status pill-steps", "The journey is the chrome — every stage visible at a glance."),
        ],
        num=9)

def s_spreading():
    _img_caption_slide(
        "Financial Spreading", "Cell provenance · override-with-reason · ratios",
        f"{ASSETS}/05-spreading.png",
        [
            ("Multi-period grid", "Years across, line items down — derived lines (grey) compute automatically."),
            ("Cell-level provenance", "Each cell shows where it came from: extracted, derived, or analyst-overridden."),
            ("Override-with-reason gate", "Editing an extracted cell past a material threshold resets confirmation."),
            ("Ratios on every period", "Leverage, ICR, DSCR, EBITDA margin — feed the rating scorecard."),
        ],
        num=10)

def s_docintel():
    _img_caption_slide(
        "Doc Intelligence", "AI extracts → human confirms",
        f"{ASSETS}/06-docintel.png",
        [
            ("AI EXTRACTS → HUMAN CONFIRMS", "Classification + structured extraction; never auto-applied."),
            ("Multilingual + translation", "Advisory — copy the text to the spread by hand."),
            ("Type-aware extraction", "PROPERTY / VEHICLE / INSURANCE / TITLE / BOND / PG — one parser per type."),
            ("Confidence routing", "Low-confidence items routed to analyst review with the source page."),
        ],
        num=11)

def s_cpt():
    _img_caption_slide(
        "Client Planning Template", "Wallet sizing · cross-sell · completeness nudges",
        f"{ASSETS}/07-cpt.png",
        [
            ("AI DRAFTS → RM CONFIRMS", "The plan is auto-assembled; the RM signs the version."),
            ("3-scenario wallet sizing", "Best / Most-Likely / Stretch over 3 years — base off confirmed EBITDA."),
            ("Cross-sell whitespace", "Catalogue diff against current facilities — heuristic, advisory."),
            ("Member figures unchanged", "Rating, capital and pricing quoted verbatim — never re-derived here."),
        ],
        num=12)

def s_groups():
    _img_caption_slide(
        "Borrower Groups", "Advisory group ID · combined credit proposal",
        f"{ASSETS}/08-groups.png",
        [
            ("AI suggests, RM tags", "Name + identifier + RM + sector fuzziness; the human still tags."),
            ("Recommendation ladder", "TAG_TO_EXISTING / REVIEW_CANDIDATES / CREATE_NEW / NO_STRONG_MATCH."),
            ("Combined credit proposal", "Insights rolled up across members; per-member figures unchanged."),
            ("Sibling discovery", "Finds ungrouped siblings of a tagged parent."),
        ],
        num=13)

def s_cad_monitoring():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "CAD · Documentation · Monitoring", crumb="Post-approval lifecycle · SoD-enforced · DMS-fed", page_num=14, page_of=PAGE_OF)
    # left card
    cw = Inches(6.0); ch = Inches(5.5); cy = Inches(1.10); gap = Inches(0.2)
    x1 = Inches(0.6); x2 = x1 + cw + gap
    add_round_rect(s, x1, cy, cw, ch, CANVAS, radius_frac=0.03, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x1+Inches(0.3), cy+Inches(0.3), cw-Inches(0.6), Inches(0.4),
             "CREDIT ADMINISTRATION (CAD)", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    add_text(s, x1+Inches(0.3), cy+Inches(0.65), cw-Inches(0.6), Inches(0.6),
             "Post-approval, before drawdown",
             size=18, color=INK, letter_spacing=-0.3)
    items_l = [
        ("Checklist from master", "Per-product checklist materialised at case open."),
        ("2-level SoD waivers", "Maker ≠ checker; L1 ≠ L2; deviations are tracked."),
        ("Limit-release trigger", "When the checklist clears, limits are released — auditably."),
        ("DMS feed", "Submitted evidence flows back into the case timeline."),
    ]
    yy = cy + Inches(1.6)
    for t, b in items_l:
        add_text(s, x1+Inches(0.3), yy, Inches(0.2), Inches(0.3), "●", size=11, color=OCEAN_5, bold=True)
        add_text(s, x1+Inches(0.55), yy, cw-Inches(0.85), Inches(0.32), t, size=12, bold=True, color=INK)
        add_text(s, x1+Inches(0.55), yy+Inches(0.32), cw-Inches(0.85), Inches(0.6), b, size=10.5, color=BODY, line_spacing=1.45)
        yy += Inches(0.85)

    add_round_rect(s, x2, cy, cw, ch, CANVAS, radius_frac=0.03, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x2+Inches(0.3), cy+Inches(0.3), cw-Inches(0.6), Inches(0.4),
             "MONITORING OF EXCEPTIONS & RENEWALS (MER)", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    add_text(s, x2+Inches(0.3), cy+Inches(0.65), cw-Inches(0.6), Inches(0.6),
             "Live register, swept daily",
             size=18, color=INK, letter_spacing=-0.3)
    items_r = [
        ("Deferred docs · conditions subsequent", "Tracked from case completion; ageing into OVERDUE."),
        ("Renewals: insurance · valuation · annual review", "Recurring obligations with reminders."),
        ("Escalation sweep", "Past-due items move OPEN → OVERDUE → ESCALATED, audited."),
        ("Maker-checker verify; waiver SoD", "Verifier ≠ submitter; waiver ≠ owner."),
    ]
    yy = cy + Inches(1.6)
    for t, b in items_r:
        add_text(s, x2+Inches(0.3), yy, Inches(0.2), Inches(0.3), "●", size=11, color=OCEAN_5, bold=True)
        add_text(s, x2+Inches(0.55), yy, cw-Inches(0.85), Inches(0.32), t, size=12, bold=True, color=INK)
        add_text(s, x2+Inches(0.55), yy+Inches(0.32), cw-Inches(0.85), Inches(0.6), b, size=10.5, color=BODY, line_spacing=1.45)
        yy += Inches(0.85)
    add_footer(s)

def s_limits_eod():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Limits, exposure & EOD", crumb="Multi-level tree · fungibility · FX revaluation · ledger reconciliation", page_num=15, page_of=PAGE_OF)
    # left half: text bullets; right half: customer-360 image
    add_text(s, Inches(0.6), Inches(1.2), Inches(6.5), Inches(0.6),
             "Limits as the contract between origination and the book",
             size=20, color=INK, letter_spacing=-0.3, line_spacing=1.2)
    items = [
        ("Multi-level limit tree",
         "Built from the approved deal — facilities, sublimits, interchangeability groups."),
        ("Product-processor APIs",
         "View / Validation / Utilisation: UTILISE · RELEASE · RESERVE · REVERSAL · override · freeze."),
        ("Country & department limits",
         "Composable caps; FI transaction workflow with maker-checker."),
        ("EOD batch — idempotent",
         "FX refresh + sanctioned-amount revaluation + utilisation reconciliation. Variances → ops desk."),
        ("Append-only utilisation ledger",
         "Every action is a row, attributed to a named user."),
    ]
    yy = Inches(2.1)
    for t, b in items:
        add_text(s, Inches(0.6), yy, Inches(0.25), Inches(0.32), "●", size=14, color=OCEAN_5, bold=True)
        add_text(s, Inches(0.9), yy, Inches(6.5), Inches(0.32), t, size=13, bold=True, color=INK)
        add_text(s, Inches(0.9), yy+Inches(0.32), Inches(6.5), Inches(0.65), b, size=11, color=BODY, line_spacing=1.5)
        yy += Inches(0.92)

    add_image(s, f"{ASSETS}/09-customer360.png", Inches(7.5), Inches(1.10), w=Inches(5.4))
    add_text(s, Inches(7.5), Inches(6.65), Inches(5.4), Inches(0.4),
             "Customer-360 — single-borrower view across services", size=10, color=MUTED, italic=True)
    add_footer(s)

def s_ai_catalogue():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "AI capabilities — by lifecycle phase", crumb="Each is advisory; each persists separately; each has a human gate", page_num=16, page_of=PAGE_OF)
    # five columns, one per phase
    phases = [
        ("Counterparty",
         ["Group identification (parent + siblings)",
          "Screening rationale (advisory)",
          "Negative-list / dedup hints"]),
        ("Origination",
         ["Document classification + extraction",
          "Type-aware collateral extraction",
          "Statement language normalisation + translation"]),
        ("Risk",
         ["Statistical RAG band",
          "Macro directional impact",
          "Goal-seek pricing optimiser"]),
        ("Decision",
         ["AI commentary (grounded)",
          "Covenant intelligence (extract + assess)",
          "CPT (wallet sizing · nudges)",
          "Doc generation + clause surgery"]),
        ("Portfolio",
         ["EWS triggers (advisory)",
          "Industry / sector outlooks",
          "Copilot (read-only fan-out)"]),
    ]
    cols = len(phases); gap = Inches(0.15); top = Inches(1.2); h = Inches(5.0)
    avail = SW - Inches(1.2) - gap*(cols-1)
    cw = avail / cols
    x = Inches(0.6)
    for head, items in phases:
        add_round_rect(s, x, top, cw, h, LIGHT_BLUE_SOFT, radius_frac=0.04, line=OCEAN_4, line_w=Pt(0.5))
        # header bar
        add_round_rect(s, x, top, cw, Inches(0.55), OCEAN_8, radius_frac=0.04)
        add_text(s, x, top, cw, Inches(0.55), head, size=12, bold=True, color=CANVAS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=0.4)
        yy = top + Inches(0.75)
        for it in items:
            # purple AI tag chip
            add_pill(s, x+Inches(0.2), yy, Inches(0.5), Inches(0.22), "AI", AI_PURPLE_SOFT, AI_PURPLE, size=8)
            add_text(s, x+Inches(0.78), yy-Inches(0.02), cw-Inches(1.0), Inches(0.6), it, size=10, color=INK, line_spacing=1.45)
            yy += Inches(0.55)
        x += cw + gap

    add_text(s, Inches(0.6), top + h + Inches(0.1), SW-Inches(1.2), Inches(0.4),
             "Every entry: advisory entity · audit.ai stamp · human-gate transition · e2e \"figure unchanged\" assertion.",
             size=10.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s)

def s_governance_pattern():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "The governance invariant", crumb="How every AI feature in Helix is built — non-negotiable", page_num=17, page_of=PAGE_OF)
    # 4-step pattern as a horizontal flow
    steps = [
        ("01", "Advisory entity",
         "Persist as its own row with advisory=true. Never mutates the authoritative figure.",
         AI_PURPLE, AI_PURPLE_SOFT),
        ("02", "AI audit stamp",
         "audit.ai(\"<capability>\", \"<EVENT>\") records who, what, when, and the rule-pack in force.",
         OCEAN_8, LIGHT_BLUE_SOFT),
        ("03", "Human gate",
         "Confirm / approve / reject by a named accountable user. SoD-enforced where the regulation requires.",
         HUMAN_GREEN, HUMAN_GREEN_SFT),
        ("04", "E2E assertion",
         "Python e2e proves the authoritative figure is byte-identical before and after the AI runs.",
         OCEAN_7, LIGHT_BLUE_SOFT),
    ]
    n = len(steps); gap = Inches(0.15); top = Inches(1.6); h = Inches(4.6)
    avail = SW - Inches(1.2) - gap*(n-1)
    sw = avail / n
    x = Inches(0.6)
    for num, head, body, color, soft in steps:
        add_round_rect(s, x, top, sw, h, soft, radius_frac=0.04, line=color, line_w=Pt(0.75))
        add_pill(s, x+Inches(0.3), top+Inches(0.3), Inches(0.6), Inches(0.28), num, color, CANVAS, size=8)
        add_text(s, x+Inches(0.3), top+Inches(0.85), sw-Inches(0.6), Inches(0.6),
                 head, size=18, bold=False, color=color, letter_spacing=-0.3, line_spacing=1.1)
        add_text(s, x+Inches(0.3), top+Inches(1.7), sw-Inches(0.6), Inches(2.5),
                 body, size=11, color=INK, line_spacing=1.6)
        x += sw + gap

    add_text(s, Inches(0.6), top + h + Inches(0.3), SW-Inches(1.2), Inches(0.5),
             "Asserted in scripts/e2e_smoke.py — 334 assertions, including every \"figure unchanged\" claim on this deck.",
             size=11, color=OCEAN_8, align=PP_ALIGN.CENTER, italic=True, bold=True)
    add_footer(s)

def s_architecture():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Architecture", crumb="9 services · SQLite-per-service · Spring Cloud Gateway · helix-common", page_num=18, page_of=PAGE_OF)
    # Top: gateway pill
    add_round_rect(s, Inches(0.6), Inches(1.2), SW-Inches(1.2), Inches(0.7), OCEAN_8, radius_frac=0.3)
    add_text(s, Inches(0.6), Inches(1.2), SW-Inches(1.2), Inches(0.7),
             "Spring Cloud Gateway · :8080 · CORS · strip-prefix · React + Vite UI",
             size=13, color=CANVAS, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=0.4)

    services = [
        ("config", "8081", "Rule packs · masters", DET_OCEAN),
        ("counterparty", "8082", "KYC · CDD · UBO · groups · screening", OCEAN_8),
        ("origination", "8083", "Apps · spread · structure · doc-intel · collateral", OCEAN_8),
        ("risk", "8084", "Rating · capital · RAROC · RAG · macro · pricing", OCEAN_8),
        ("decision", "8085", "DoA · CAD · MER · covenants · docs · commentary · CPT · groups", OCEAN_8),
        ("portfolio", "8086", "Book · ECL · EWS · MIS · CAP · exports", OCEAN_8),
        ("limit", "8088", "Tree · utilisation · EOD · ledger · FI", OCEAN_8),
        ("copilot", "8087", "Read-only conversational fan-out", AI_PURPLE),
    ]
    cols = 4; gap = Inches(0.15); top = Inches(2.2); h = Inches(1.2)
    avail = SW - Inches(1.2) - gap*(cols-1)
    cw = avail / cols
    for i, (name, port, sub, c) in enumerate(services):
        r, col = divmod(i, cols)
        x = Inches(0.6) + col*(cw+gap)
        y = top + r*(h+Inches(0.18))
        add_round_rect(s, x, y, cw, h, CANVAS, radius_frac=0.06, line=HAIRLINE, line_w=Pt(0.5))
        # port chip
        add_pill(s, x+cw-Inches(0.7), y+Inches(0.15), Inches(0.55), Inches(0.24), port, LIGHT_BLUE_SOFT, OCEAN_7, size=8, letter_spacing=0)
        add_text(s, x+Inches(0.2), y+Inches(0.15), cw-Inches(1.0), Inches(0.34),
                 name, size=14, bold=True, color=c, letter_spacing=-0.2)
        add_text(s, x+Inches(0.2), y+Inches(0.55), cw-Inches(0.3), Inches(0.65),
                 sub, size=10, color=BODY, line_spacing=1.45)

    # Bottom: helix-common strip
    by = Inches(5.85)
    add_round_rect(s, Inches(0.6), by, SW-Inches(1.2), Inches(0.7), LIGHT_BLUE_SOFT, radius_frac=0.3, line=OCEAN_4, line_w=Pt(0.5))
    add_text(s, Inches(0.6), by, SW-Inches(1.2), Inches(0.7),
             "helix-common · canonical enums · append-only audit · ingest (Envelope/Connector/Guard) · export (Envelope/ErmRiskRecord/FinanceProvisionEntry/CprPortfolioLine)",
             size=11, color=OCEAN_8, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.6), Inches(6.75), SW-Inches(1.2), Inches(0.4),
             "Java 25 · Spring Boot · SQLite-per-service (Hikari pool=1, single-writer) · ddl-auto=update · Hibernate community SQLite dialect",
             size=10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s)

def s_regime_as_data():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Regime is overlay data — never code", crumb="Versioned rule packs + generic master-data engine", page_num=19, page_of=PAGE_OF)
    add_text(s, Inches(0.6), Inches(1.2), Inches(12.2), Inches(0.6),
             "A new jurisdiction is a pack, not a release.",
             size=24, color=INK, letter_spacing=-0.4, line_spacing=1.2)
    add_text(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(0.7),
             "config-service owns versioned, dual-signed rule packs + a generic master-data engine. "
             "All downstream engines fetch packs at runtime; no engine branches on country code.",
             size=12.5, color=BODY, line_spacing=1.55)

    # Two columns: pack types | master types
    cw = Inches(6.0); ch = Inches(3.5); cy = Inches(3.0); gap = Inches(0.4)
    x1 = Inches(0.6); x2 = x1 + cw + gap
    add_round_rect(s, x1, cy, cw, ch, CANVAS, radius_frac=0.04, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x1+Inches(0.3), cy+Inches(0.25), cw-Inches(0.6), Inches(0.4),
             "RULE PACKS (versioned, dual-signed)", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    packs_text = ("CAPITAL  ·  ECRA  ·  PD/LGD  ·  PROVISIONING  ·  DoA  ·  "
                  "LIMITS  ·  PRICING  ·  WORKFLOW  ·  RATING SCORECARD")
    add_text(s, x1+Inches(0.3), cy+Inches(0.7), cw-Inches(0.6), Inches(2.5),
             packs_text, size=13, color=INK, line_spacing=1.7, font=FONT_MONO)

    add_round_rect(s, x2, cy, cw, ch, CANVAS, radius_frac=0.04, line=HAIRLINE, line_w=Pt(0.75))
    add_text(s, x2+Inches(0.3), cy+Inches(0.25), cw-Inches(0.6), Inches(0.4),
             "MASTER DATA (maker-checker SoD, 22 types)", size=10, bold=True, color=OCEAN_8, letter_spacing=1.2)
    masters_text = ("DEDUP  ·  NEGATIVE LIST  ·  COVENANT LIB  ·  FACILITY MASTER  ·  "
                    "COLLATERAL MASTER  ·  RAROC  ·  EWS TRIGGERS  ·  EMAIL TPL  ·  "
                    "CHECKLISTS  ·  DOC TEMPLATES  ·  T&C  ·  …")
    add_text(s, x2+Inches(0.3), cy+Inches(0.7), cw-Inches(0.6), Inches(2.5),
             masters_text, size=12, color=INK, line_spacing=1.7, font=FONT_MONO)
    add_footer(s)

def s_audit():
    _img_caption_slide(
        "Audit Trail", "Examiner-ready · every action attributed",
        f"{ASSETS}/10-audit.png",
        [
            ("HUMAN · AI · SYSTEM", "Every row stamped with actor type. No anonymous writes."),
            ("Named accountable user", "X-Actor header travels with every write; persisted in the row."),
            ("Service-scoped + subject-scoped", "Browse by service or by subject (CP, application, deal)."),
            ("Immutable, append-only", "Hibernate ddl-only insert; rule-pack version captured per event."),
        ],
        num=20)

def s_copilot():
    _img_caption_slide(
        "Copilot", "Conversational, grounded, non-binding",
        f"{ASSETS}/11-copilot.png",
        [
            ("8 personas, scoped read-only", "RM · analyst · credit officer · committee · compliance · ops · PM · CRO."),
            ("Reads across all 9 services", "Fan-out behind a single endpoint; every answer cites its source."),
            ("Refuses credit-consequential actions", "Asks routed to the gated workflow instead of being executed."),
            ("Persona-aware tone", "CRO-speak vs analyst-speak vs operations-speak; same facts, different register."),
        ],
        num=21)

def s_exports():
    _img_caption_slide(
        "Downstream Exports", "ERM · Finance/GL · CPR · idempotent",
        f"{ASSETS}/12-exports.png",
        [
            ("Canonical export contract", "helix-common.export — typed envelopes per destination."),
            ("Idempotent by as-of day", "Re-running the same as-of returns the existing batch, not a duplicate."),
            ("Symmetric with ingest", "Same envelope discipline as inbound connectors."),
            ("Engine-stamped audit", "audit.engine(\"EXPORT_GENERATED\") on every batch."),
        ],
        num=22)

def s_what_makes_helix_different():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "What makes Helix different", crumb="Three claims, each provable", page_num=23, page_of=PAGE_OF)
    rows = [
        ("01",
         "AI is built into the boundaries, not the figure path",
         "Extraction, drafting, scoring, optimisation — all live as separate advisory entities, "
         "never as a mutation. e2e proves the figure is byte-identical before and after."),
        ("02",
         "Governance is in the chrome, not just the docs",
         "Every screen carries AI / HUMAN / DETERMINISTIC chips. The Risk Lab is literally split into "
         "advisory and authoritative panes. The Pricing Lab carries a PRESERVED tag."),
        ("03",
         "Regime is overlay data — never code branches",
         "A new jurisdiction is a versioned rule pack, not a release. Engines fetch packs at runtime. "
         "22 master types, all maker-checker, all SoD-aware."),
    ]
    top = Inches(1.3); h = Inches(1.7); gap = Inches(0.2)
    y = top
    for num, head, body in rows:
        add_round_rect(s, Inches(0.6), y, SW-Inches(1.2), h, CANVAS, radius_frac=0.04, line=HAIRLINE, line_w=Pt(0.75))
        add_pill(s, Inches(0.85), y+Inches(0.25), Inches(0.6), Inches(0.32), num, OCEAN_8, CANVAS, size=10)
        add_text(s, Inches(1.65), y+Inches(0.22), Inches(11), Inches(0.5),
                 head, size=17, bold=True, color=INK, letter_spacing=-0.3)
        add_text(s, Inches(1.65), y+Inches(0.72), Inches(11), Inches(0.95),
                 body, size=12, color=BODY, line_spacing=1.55)
        y += h + gap
    add_footer(s)

def s_metrics():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "By the numbers", crumb="Built. Tested. Reproducible.", page_num=24, page_of=PAGE_OF)
    stats = [
        ("9", "microservices", "Spring Boot, SQLite-per-service"),
        ("22", "master types", "maker-checker, SoD-aware"),
        ("9", "rule-pack kinds", "versioned, dual-signed"),
        ("334", "e2e assertions", "Python driver, gateway-routed"),
        ("100", "obligor stress run", "distributed book, ~60s"),
        ("8", "copilot personas", "scoped, read-only"),
        ("AI · HUMAN · SYSTEM", "actor types", "stamped on every audit row"),
        ("Java 25", "stack", "Spring · React · Vite · TS"),
    ]
    cols = 4; rows = 2; gap = Inches(0.15); top = Inches(1.3); ch = Inches(2.5)
    avail = SW - Inches(1.2) - gap*(cols-1)
    cw = avail / cols
    for i, (v, lbl, sub) in enumerate(stats):
        r, c = divmod(i, cols)
        x = Inches(0.6) + c*(cw+gap)
        y = top + r*(ch+Inches(0.2))
        add_round_rect(s, x, y, cw, ch, LIGHT_BLUE_SOFT, radius_frac=0.04, line=OCEAN_4, line_w=Pt(0.5))
        # number
        add_text(s, x+Inches(0.2), y+Inches(0.3), cw-Inches(0.4), Inches(0.95),
                 v, size=30, bold=True, color=OCEAN_8, font=FONT_MONO, align=PP_ALIGN.CENTER, letter_spacing=-0.5, line_spacing=1.0)
        add_text(s, x+Inches(0.2), y+Inches(1.3), cw-Inches(0.4), Inches(0.4),
                 lbl, size=12, bold=True, color=OCEAN_7, align=PP_ALIGN.CENTER, letter_spacing=0.4)
        add_text(s, x+Inches(0.2), y+Inches(1.75), cw-Inches(0.4), Inches(0.65),
                 sub, size=10, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.45, italic=True)
    add_footer(s)

def s_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_topbar(s, "Roadmap", crumb="Each beat extends governed AI further into the workflow", page_num=25, page_of=PAGE_OF)
    rows = [
        ("Q3", "Statistical EWS thresholds + borrower-level RAG scoring",
         "Move from rules-based EWS to a model-fit-aware band. Borrower RAG sits beside the deal RAG."),
        ("Q3", "Credit-history & news ingestion connector",
         "Canonical connector for credit-bureau + news; advisory enrichment, human-gated tagging."),
        ("Q4", "Portfolio-360 AI commentary + regulatory narrative",
         "Auto-drafted board narrative grounded on the same figures — confirm-locked before release."),
        ("Q4", "Exception highlights pane",
         "Top-of-screen drawer surfacing the actionable exceptions across the user's portfolio."),
        ("H1 next", "Multi-jurisdiction simultaneously",
         "Two rule-pack sets live in one tenancy; deals tagged to a jurisdiction; same code."),
    ]
    top = Inches(1.3); h = Inches(1.0); gap = Inches(0.15)
    y = top
    for q, head, body in rows:
        add_round_rect(s, Inches(0.6), y, SW-Inches(1.2), h, CANVAS, radius_frac=0.06, line=HAIRLINE, line_w=Pt(0.5))
        add_pill(s, Inches(0.85), y+Inches(0.30), Inches(0.85), Inches(0.34), q, OCEAN_8, CANVAS, size=10, letter_spacing=0.4)
        add_text(s, Inches(1.95), y+Inches(0.20), Inches(10.5), Inches(0.4),
                 head, size=14, bold=True, color=INK, letter_spacing=-0.2)
        add_text(s, Inches(1.95), y+Inches(0.55), Inches(10.5), Inches(0.5),
                 body, size=11, color=BODY, line_spacing=1.5)
        y += h + gap
    add_footer(s)

def s_close():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, SURFACE_DARK)
    # accent
    accent = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-1), Inches(9), Inches(9))
    accent.fill.solid(); accent.fill.fore_color.rgb = OCEAN_6; accent.line.fill.background()
    _disable_shadow(accent)
    sppr = accent.fill._xPr.find(qn('a:solidFill'))
    sppr.find(qn('a:srgbClr')).append(etree.fromstring('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="14000"/>'))

    add_text(s, Inches(0.85), Inches(1.6), Inches(11), Inches(1.2),
             "AI you can defend to regulators.",
             size=48, color=CANVAS, letter_spacing=-1.2, line_spacing=1.05)
    add_text(s, Inches(0.85), Inches(3.0), Inches(11), Inches(0.6),
             "Governed AI · Deterministic figures · End-to-end wholesale credit.",
             size=20, color=RGBColor(0x9F, 0xB8, 0xC0), letter_spacing=-0.2)

    # CTA card
    add_round_rect(s, Inches(0.85), Inches(4.2), Inches(11.6), Inches(2.0), SURFACE_DARK_2, radius_frac=0.04, line=OCEAN_6, line_w=Pt(0.5))
    add_text(s, Inches(1.1), Inches(4.4), Inches(11), Inches(0.4),
             "TAKE THE NEXT STEP", size=10, bold=True, color=OCEAN_5, letter_spacing=1.5)
    add_text(s, Inches(1.1), Inches(4.8), Inches(11), Inches(0.5),
             "Schedule a 30-minute walk-through against your portfolio's data.",
             size=18, color=CANVAS, letter_spacing=-0.3)
    add_text(s, Inches(1.1), Inches(5.5), Inches(11), Inches(0.6),
             "Live demo · governance invariants asserted in your jurisdiction's rule pack · sample data masked.",
             size=12, color=RGBColor(0x9F, 0xB8, 0xC0), line_spacing=1.5)

    add_text(s, Inches(0.85), Inches(6.8), Inches(11), Inches(0.4),
             "CRISIL  ·  Helix — Governed AI for Wholesale Credit",
             size=10, color=RGBColor(0x9F, 0xB8, 0xC0), letter_spacing=1.5)

# ─── Assemble ───────────────────────────────────────────────────────────────
s_title()              # 1
s_problem()            # 2
s_thesis()             # 3
s_what_is_helix()      # 4
s_lifecycle()          # 5
s_dashboard()          # 6
s_risk_lab_split()     # 7
s_pricing_lab()        # 8
s_workspace()          # 9
s_spreading()          # 10
s_docintel()           # 11
s_cpt()                # 12
s_groups()             # 13
s_cad_monitoring()     # 14
s_limits_eod()         # 15
s_ai_catalogue()       # 16
s_governance_pattern() # 17
s_architecture()       # 18
s_regime_as_data()     # 19
s_audit()              # 20
s_copilot()            # 21
s_exports()            # 22
s_what_makes_helix_different()  # 23
s_metrics()            # 24
s_roadmap()            # 25
s_close()              # 26

out = "/home/user/CPS-LOS/docs/Helix-Pitch-Deck.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides)} slides)")
